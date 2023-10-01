"""
Convert Twine Proof file to PowerPoint file

Usage:
  convertTwineProofToPowerPoint.py <in_file> <out_file>
  convertTwineProofToPowerPoint.py -h | --help
  convertTwineProofToPowerPoint.py --version

  <in_file>: filename of the Twine file to be converted
  <out_file>: output filename of PowerPoint

Examples:
  convertTwineProofToPowerPoint.py examples/sample.html converted.pptx
    "sample.html" in the folder is converted to a PPTX file
    under the filename of "converted.pptx".

Options:
  -h --help     Show this screen.
  --version     Show version.
"""

import re
from docopt import docopt
from pptx import Presentation
from pptx.util import Inches, Pt

def my_index_multi(l, x):
    return [i for i, _x in enumerate(l) if _x[0] == x]

if __name__ == '__main__':

    arguments = docopt(__doc__, version="0.1")
    in_file = arguments["<in_file>"]
    out_file = arguments["<out_file>"]

    xml_original_filename = in_file
    with open(xml_original_filename, "r",encoding="utf-8") as f:
        xml = f.read()
    pattern = re.compile(r'<tw-passagedata.*?>.*?</tw-passagedata>', re.MULTILINE | re.DOTALL)
    inside_text = re.compile(r'<tw-passagedata.*?>(.*?)</tw-passagedata>', re.MULTILINE | re.DOTALL)
    inside_hyperlink = re.compile(r'\[\[(.*?)\]\]', re.MULTILINE | re.DOTALL)

    slides = []
    links = []

    result = pattern.finditer( xml)

    for m in result:
        print("-->",m.group())
        result_name = re.finditer(r'<tw-passagedata.*?name="(.*?)".*?>',m.group())
        inside_text_string = inside_text.finditer(m.group())
        for name in result_name:
            # print(name.groups()[0])
            header_text = name.groups()[0]
            for in_text in inside_text_string:
                print(in_text.groups()[0])
                inside_text_parsed = in_text.groups()[0]
                slides.append([header_text,inside_text_parsed])
                found_links = inside_hyperlink.finditer(inside_text_parsed)
                for each_link in found_links:
                    each_link_string = each_link.groups()[0]
                    print(each_link_string)
                    links.append([header_text,each_link_string])

    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    for a, b in slides:
        slide = prs.slides.add_slide(blank_slide_layout)

        left = top = width = height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        if len(b) > 20000:
            tf.text = a + "\n" + b[0:200]
            print("cut off at 200",b[0:200])
            print("removed at and after 201", b[201:])
        else:
            tf.text = a + "\n" + b

    left = Inches(5)
    for i, (a, b) in enumerate(links):
        found_slide_numbers = my_index_multi(slides,b)
        ref_item = None
        if len(found_slide_numbers) != 0:
            ref_item = found_slide_numbers
        else:
            continue

        for i_slide, item in enumerate(slides):
            if item[0] == a:
                len_of_shapes = len(prs.slides[i_slide].shapes)
                the_element_order = len_of_shapes * 0.2
                top = Inches(the_element_order)
                txBox = prs.slides[i_slide].shapes.add_textbox(left, top, width, height)
                p = txBox.text_frame.add_paragraph()
                p.text = b
                click_action = txBox.click_action
                if len(ref_item) == 1:
                    click_action.target_slide = prs.slides[ref_item[0]]
                else:
                    print("ref_item more than 1",ref_item)

    prs.save(out_file)